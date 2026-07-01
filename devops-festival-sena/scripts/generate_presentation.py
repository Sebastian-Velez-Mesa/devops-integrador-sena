import os
import sys
import subprocess

# ── Auto-instalación de dependencias ──────────────────────────────────
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("La biblioteca 'python-pptx' no está instalada. Instalándola ahora...")
    try:
        subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        print("Biblioteca instalada y cargada exitosamente.")
    except Exception as e:
        print(f"Error al instalar python-pptx: {e}")
        print("Asegúrate de tener conexión a internet y pip instalado.")
        sys.exit(1)

# ── Paleta de Colores Corporativos (Alineados con el sitio web) ───────
COLOR_BG = RGBColor(11, 13, 25)       # HSL oscuro profundo
COLOR_TEXT_MAIN = RGBColor(240, 240, 245) # Blanco suave
COLOR_MUTED = RGBColor(180, 180, 200) # Gris suave
COLOR_PRIMARY = RGBColor(168, 85, 247) # Púrpura brillante
COLOR_SECONDARY = RGBColor(6, 182, 212) # Cian brillante

def set_slide_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def create_text_box(slide, left, top, width, height):
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)
    return tf

def add_title(slide, text):
    tf = create_text_box(slide, Inches(0.75), Inches(0.6), Inches(11.83), Inches(1.2))
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Outfit"
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = COLOR_SECONDARY
    return tf

def main():
    prs = Presentation()
    
    # Configurar formato Widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6] # Layout en blanco

    # ──────────────────────────────────────────────────────────────────
    # SLIDE 1: Portada (Slide de Título)
    # ──────────────────────────────────────────────────────────────────
    slide_1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide_1, COLOR_BG)
    
    tf_1 = create_text_box(slide_1, Inches(0.75), Inches(1.8), Inches(11.83), Inches(4.5))
    
    p_badge = tf_1.paragraphs[0]
    p_badge.text = "SITUACIÓN DE APRENDIZAJE INTEGRADORA - DEVOPS"
    p_badge.font.name = "Inter"
    p_badge.font.size = Pt(14)
    p_badge.font.bold = True
    p_badge.font.color.rgb = COLOR_PRIMARY
    p_badge.space_after = Pt(15)
    
    p_title = tf_1.add_paragraph()
    p_title.text = "SENA DevOps & Cloud Fest 2026"
    p_title.font.name = "Outfit"
    p_title.font.size = Pt(54)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_TEXT_MAIN
    p_title.space_after = Pt(8)

    p_sub = tf_1.add_paragraph()
    p_sub.text = "Sustentación Técnica del Proyecto Multicontenedor y Pipeline CI/CD"
    p_sub.font.name = "Inter"
    p_sub.font.size = Pt(22)
    p_sub.font.color.rgb = COLOR_SECONDARY
    p_sub.space_after = Pt(45)

    p_info = tf_1.add_paragraph()
    p_info.text = "Programa de Formación: DevOps y Contenedores (Docker)\nCentro de Tecnología y Manufactura Avanzada (CTMA)\nInstructor: Elkin Moreno"
    p_info.font.name = "Inter"
    p_info.font.size = Pt(13)
    p_info.font.color.rgb = COLOR_MUTED

    # ──────────────────────────────────────────────────────────────────
    # SLIDE 2: Introducción y Contextualización
    # ──────────────────────────────────────────────────────────────────
    slide_2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide_2, COLOR_BG)
    add_title(slide_2, "1. Introducción y Contextualización")
    
    tf_2 = create_text_box(slide_2, Inches(0.75), Inches(2.0), Inches(11.83), Inches(4.8))
    
    bullets_2 = [
        ("Nombre del Proyecto:", "SENA DevOps & Cloud Fest 2026 - Portal Web Académico/Cultural."),
        ("Problema que Resuelve:", "Falta de demostración práctica en vivo de infraestructuras DevOps. Este proyecto unifica desarrollo ágil (Git Flow), contenedores e integración continua en una interfaz viva que expone sus propias métricas."),
        ("Objetivo General:", "Integrar y demostrar los conocimientos adquiridos en las semanas 1 a 7 (Docker, Redes, Volúmenes, Git, Pipelines CI/CD) mediante una aplicación multicontenedor altamente visual."),
        ("Usuarios Beneficiados:", "Comité evaluador, instructores y aprendices del programa de tecnología, sirviendo de plantilla base de producción para proyectos escolares."),
        ("Tecnologías Core:", "Frontend Nginx (Alpine), Backend Flask (Python 3.11), WSGI Gunicorn, GitHub Actions, Docker & Docker Compose.")
    ]
    
    for i, (title, text) in enumerate(bullets_2):
        p = tf_2.paragraphs[0] if i == 0 else tf_2.add_paragraph()
        p.text = "• "
        p.space_after = Pt(16)
        
        run_title = p.add_run()
        run_title.text = title + " "
        run_title.font.name = "Inter"
        run_title.font.size = Pt(16)
        run_title.font.bold = True
        run_title.font.color.rgb = COLOR_PRIMARY
        
        run_text = p.add_run()
        run_text.text = text
        run_text.font.name = "Inter"
        run_text.font.size = Pt(16)
        run_text.font.color.rgb = COLOR_TEXT_MAIN

    # ──────────────────────────────────────────────────────────────────
    # SLIDE 3: Arquitectura General
    # ──────────────────────────────────────────────────────────────────
    slide_3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide_3, COLOR_BG)
    add_title(slide_3, "2. Arquitectura del Proyecto")
    
    tf_3 = create_text_box(slide_3, Inches(0.75), Inches(2.0), Inches(11.83), Inches(4.8))
    
    # Descripción de la topología
    p_arch_desc = tf_3.paragraphs[0]
    p_arch_desc.text = "La solución implementa una topología desacoplada en base a microservicios orquestados:"
    p_arch_desc.font.name = "Inter"
    p_arch_desc.font.size = Pt(18)
    p_arch_desc.font.bold = True
    p_arch_desc.font.color.rgb = COLOR_TEXT_MAIN
    p_arch_desc.space_after = Pt(20)
    
    bullets_3 = [
        ("Frontend Nginx (puerto 8080:80):", "Sirve el código estático (HTML, CSS, JS). Realiza peticiones asíncronas desde el navegador del cliente hacia el backend."),
        ("Backend Flask (puerto 5000:5000):", "API REST de Python que procesa solicitudes del festival y métricas de infraestructura. Servido por Gunicorn."),
        ("Red Bridge Privada ('devops-network'):", "Garantiza aislamiento. Los contenedores se descubren por nombre DNS interno sin necesidad de mapear puertos del host públicamente."),
        ("Volumen Persistente ('devops-backend-logs'):", "Almacena los logs (/app/logs/festival.log) y registros de contacto (/app/logs/contactos.json). Los datos persisten aunque se re-inicialice el contenedor backend.")
    ]
    
    for title, text in bullets_3:
        p = tf_3.add_paragraph()
        p.text = "➤ "
        p.space_after = Pt(12)
        
        run_title = p.add_run()
        run_title.text = title + " "
        run_title.font.name = "Inter"
        run_title.font.size = Pt(15)
        run_title.font.bold = True
        run_title.font.color.rgb = COLOR_SECONDARY
        
        run_text = p.add_run()
        run_text.text = text
        run_text.font.name = "Inter"
        run_text.font.size = Pt(15)
        run_text.font.color.rgb = COLOR_TEXT_MAIN

    # ──────────────────────────────────────────────────────────────────
    # SLIDE 4: Decisiones Técnicas y Dockerización
    # ──────────────────────────────────────────────────────────────────
    slide_4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide_4, COLOR_BG)
    add_title(slide_4, "3. Decisiones de Diseño y Contenedores")
    
    tf_4 = create_text_box(slide_4, Inches(0.75), Inches(2.0), Inches(11.83), Inches(4.8))
    
    bullets_4 = [
        ("Optimización de Imagen Base:", "Utilización de 'python:3.11-slim' para minimizar el tamaño de la imagen (~150MB en vez de ~900MB de Python completo), reduciendo superficie de ataque y agilizando despliegues."),
        ("Seguridad (Ejecución sin privilegios):", "Configuración de un usuario y grupo del sistema sin privilegios ('devopsuser') en el Dockerfile del backend. Si el contenedor se ve comprometido, el atacante no obtendrá acceso root."),
        ("Comprobación de Salud Activa (Healthcheck):", "El backend define un comando healthcheck nativo usando curl en el puerto 5000/health, permitiendo que Compose determine si el servicio está operativo."),
        ("Dependencias de Inicio Condicionales (depends_on):", "El frontend Nginx arranca únicamente después de que el backend pasa satisfactoriamente su control de salud, evitando peticiones vacías al inicio."),
        ("Capa de Logs y Persistencia:", "Los logs son formateados a nivel de aplicación e interactúan directamente con el Dashboard DevOps en el frontend.")
    ]
    
    for i, (title, text) in enumerate(bullets_4):
        p = tf_4.paragraphs[0] if i == 0 else tf_4.add_paragraph()
        p.text = "✔ "
        p.space_after = Pt(14)
        
        run_title = p.add_run()
        run_title.text = title + " "
        run_title.font.name = "Inter"
        run_title.font.size = Pt(15)
        run_title.font.bold = True
        run_title.font.color.rgb = COLOR_PRIMARY
        
        run_text = p.add_run()
        run_text.text = text
        run_text.font.name = "Inter"
        run_text.font.size = Pt(15)
        run_text.font.color.rgb = COLOR_TEXT_MAIN

    # ──────────────────────────────────────────────────────────────────
    # SLIDE 5: Automatización CI/CD (GitHub Actions)
    # ──────────────────────────────────────────────────────────────────
    slide_5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide_5, COLOR_BG)
    add_title(slide_5, "4. Pipeline de CI/CD Automatizado")
    
    tf_5 = create_text_box(slide_5, Inches(0.75), Inches(2.0), Inches(11.83), Inches(4.8))
    
    p_cicd_desc = tf_5.paragraphs[0]
    p_cicd_desc.text = "GitHub Actions Workflow (.github/workflows/ci-cd.yml) para automatización completa:"
    p_cicd_desc.font.name = "Inter"
    p_cicd_desc.font.size = Pt(18)
    p_cicd_desc.font.bold = True
    p_cicd_desc.font.color.rgb = COLOR_TEXT_MAIN
    p_cicd_desc.space_after = Pt(20)
    
    bullets_5 = [
        ("Disparadores (Triggers):", "Ejecutado automáticamente ante cualquier Push o Pull Request en las ramas 'main' y 'develop'."),
        ("Fase de Control de Calidad (Linting):", "Usa 'flake8' para auditar errores de sintaxis y guías de estilo en Python, previniendo fallos en producción."),
        ("Fase de Pruebas Unitarias (Testing):", "Ejecuta 'pytest' sobre el backend para validar que todas las rutas clave de la API (/, /health, /api/artistas, /api/contacto) devuelvan estados HTTP válidos."),
        ("Fase de Validación de Compilación (Build):", "Realiza builds simulados de las imágenes Docker (backend y frontend) para asegurar que el código compile y los archivos Dockerfile sean correctos antes de desplegar.")
    ]
    
    for title, text in bullets_5:
        p = tf_5.add_paragraph()
        p.text = "⚡ "
        p.space_after = Pt(12)
        
        run_title = p.add_run()
        run_title.text = title + " "
        run_title.font.name = "Inter"
        run_title.font.size = Pt(15)
        run_title.font.bold = True
        run_title.font.color.rgb = COLOR_SECONDARY
        
        run_text = p.add_run()
        run_text.text = text
        run_text.font.name = "Inter"
        run_text.font.size = Pt(15)
        run_text.font.color.rgb = COLOR_TEXT_MAIN

    # ──────────────────────────────────────────────────────────────────
    # SLIDE 6: Conclusiones
    # ──────────────────────────────────────────────────────────────────
    slide_6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide_6, COLOR_BG)
    add_title(slide_6, "5. Conclusiones y Aprendizajes")
    
    tf_6 = create_text_box(slide_6, Inches(0.75), Inches(2.0), Inches(11.83), Inches(4.8))
    
    bullets_6 = [
        ("Dificultad Encontrada / Solución:", "Establecer la comunicación inter-contenedores. Se resolvió estructurando una red bridge de Docker y vinculando las consultas fetch a la IP pública del host de forma asíncrona."),
        ("Cultura DevOps Aplicada:", "Automatización desde el primer commit. El testing e integración continua garantizan la estabilidad del software eliminando el clásico 'en mi máquina funciona'."),
        ("Aprendizaje Adquirido:", "Comprensión profunda de la gestión de redes Docker, montajes de volúmenes persistentes y encapsulamiento de dependencias de sistema mediante Dockerfiles optimizados."),
        ("Mejoras Futuras:", "Implementación de una base de datos real persistente (e.g., PostgreSQL) como contenedor adicional en Compose, encriptación TLS con Nginx proxy reverso y orquestación con Kubernetes.")
    ]
    
    for i, (title, text) in enumerate(bullets_6):
        p = tf_6.paragraphs[0] if i == 0 else tf_6.add_paragraph()
        p.text = "💡 "
        p.space_after = Pt(15)
        
        run_title = p.add_run()
        run_title.text = title + " "
        run_title.font.name = "Inter"
        run_title.font.size = Pt(15)
        run_title.font.bold = True
        run_title.font.color.rgb = COLOR_PRIMARY
        
        run_text = p.add_run()
        run_text.text = text
        run_text.font.name = "Inter"
        run_text.font.size = Pt(15)
        run_text.font.color.rgb = COLOR_TEXT_MAIN

    # Guardar presentación en el directorio raíz del proyecto
    filename = "Sustentacion_DevOps.pptx"
    out_path = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), filename)
    prs.save(out_path)
    print(f"Presentación de PowerPoint generada exitosamente en: {out_path}")

if __name__ == "__main__":
    main()
